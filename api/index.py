import os
import json
import urllib.parse
import mimetypes
from http.server import HTTPServer, BaseHTTPRequestHandler
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml import parse_xml, OxmlElement
from docx.oxml.ns import nsdecls, qn

CURRENT_DIR = os.path.dirname(os.path.abspath(__file__))
DOWNLOAD_DIR = os.path.join(CURRENT_DIR, 'static')
os.makedirs(DOWNLOAD_DIR, exist_ok=True)

PPT_TEMPLATE_PATH = os.path.join(CURRENT_DIR, 'template.pptx')
WORD_TEMPLATE_PATH = os.path.join(CURRENT_DIR, 'template.docx')

# ---------------------------------------------------------------------------
# 1. PPT 生成逻辑
import os
import json
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER

# 1. 基准路径解析
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PPT_TEMPLATE_PATH = os.path.join(BASE_DIR, "template.pptx")

def safe_json_parse(data):
    """安全解析字典或 JSON 字符串，兼容嵌套转义与顶层包裹"""
    if isinstance(data, dict):
        return data
    if isinstance(data, str):
        try:
            parsed = json.loads(data)
            if isinstance(parsed, str):  # 防止二次转义的 JSON 字符串
                parsed = json.loads(parsed)
            return parsed
        except Exception:
            return {}
    return {}

def extract_llm_payload(raw_data):
    """
    智能解包 LLM 的输出：自动识别 structured_output、text、word_data 等不同层级，
    确保提取出完整的 basic_info 与 word_data。
    """
    data = safe_json_parse(raw_data)
    
    # 优先提取 structured_output
    if "structured_output" in data and isinstance(data["structured_output"], dict):
        data = data["structured_output"]

    basic_info = data.get("basic_info", {})
    
    # 获取核心数据内容（兼容 word_data / ppt_content / 根节点）
    word_data = data.get("word_data") or data.get("ppt_content") or data

    # 防护：如果 word_data 包含 basic_info，再次合并
    if isinstance(word_data, dict) and "basic_info" in word_data:
        if not basic_info:
            basic_info = word_data.get("basic_info", {})

    return basic_info, word_data

def get_layout_by_type(prs):
    """
    智能挑选母版布局：
    优先按名称匹配 (Cover/Title/封面 与 Content/Body/正文)，兜底按索引 0 和 1
    """
    cover_layout = None
    content_layout = None

    for layout in prs.slide_layouts:
        name = layout.name.lower()
        if any(k in name for k in ["cover", "title", "封面", "首页"]) and not cover_layout:
            cover_layout = layout
        elif any(k in name for k in ["content", "body", "正文", "内容", "内页"]) and not content_layout:
            content_layout = layout

    # 兜底策略
    if not cover_layout and len(prs.slide_layouts) > 0:
        cover_layout = prs.slide_layouts[0]
    if not content_layout:
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else cover_layout

    return cover_layout, content_layout

def find_placeholders(slide):
    """
    精准识别页面中的占位符：标题、副标题、正文内容
    """
    title_ph = slide.shapes.title
    subtitle_ph = None
    body_ph = None

    for shape in slide.placeholders:
        if not shape.has_text_frame:
            continue
        ph_type = shape.placeholder_format.type
        
        # 匹配标题
        if ph_type in (PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE) and not title_ph:
            title_ph = shape
        # 匹配副标题
        elif ph_type == PP_PLACEHOLDER.SUBTITLE and not subtitle_ph:
            subtitle_ph = shape
        # 匹配正文内容/对象
        elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) and not body_ph:
            body_ph = shape

    # 若未按类型识别出，则按索引或名称进行二次匹配
    if not subtitle_ph or not body_ph:
        for shape in slide.placeholders:
            if shape == title_ph or not shape.has_text_frame:
                continue
            name = shape.name.lower()
            if ("subtitle" in name or shape.placeholder_format.idx == 1) and not subtitle_ph:
                subtitle_ph = shape
            elif ("content" in name or "body" in name or "text" in name) and not body_ph:
                body_ph = shape

    return title_ph, subtitle_ph, body_ph

def fill_text_frame(tf, text_content, default_size=None, align=None, default_color=None):
    """
    填充文本到 TextFrame，保留母版原本样式：
    1. 清空默认 Placeholder 示例文本
    2. 支持按 \n 拆分成多个独立段落
    3. 自动识别 【标签】 并加粗处理
    """
    tf.word_wrap = True
    tf.text = ""  # 清空占位符里的默认文本

    if isinstance(text_content, str):
        lines = [line.strip() for line in text_content.split("\n") if line.strip()]
    elif isinstance(text_content, list):
        lines = text_content
    else:
        lines = [str(text_content)]

    for idx, line_text in enumerate(lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        if align is not None:
            p.alignment = align

        # 识别带括号的标题项（例如：【教学重点】、【教师活动】）
        if line_text.startswith("【") and "】" in line_text:
            split_idx = line_text.index("】") + 1
            header_str = line_text[:split_idx]
            body_str = line_text[split_idx:]

            # 标题部分：加粗
            run_h = p.add_run()
            run_h.text = header_str
            run_h.font.bold = True
            if default_size: run_h.font.size = Pt(default_size)
            if default_color: run_h.font.color.rgb = RGBColor(*default_color)

            # 正文部分：常态
            if body_str:
                run_b = p.add_run()
                run_b.text = body_str
                if default_size: run_b.font.size = Pt(default_size)
                if default_color: run_b.font.color.rgb = RGBColor(*default_color)
        else:
            run = p.add_run()
            run.text = line_text
            if default_size: run.font.size = Pt(default_size)
            if default_color: run.font.color.rgb = RGBColor(*default_color)

def generate_ppt_from_template(ppt_data, basic_info=None, output_path="output.pptx"):
    try:
        # 1. 动态抽取与解析 LLM 返回的数据结构
        parsed_basic_info, word_data = extract_llm_payload(ppt_data)
        
        # 补充外部传入的 basic_info (如果有)
        ext_basic = safe_json_parse(basic_info) if basic_info else {}
        final_basic_info = {**parsed_basic_info, **ext_basic}

        # 2. 读取 PPT 模板
        template_exists = os.path.exists(PPT_TEMPLATE_PATH)
        print(f"[PPT Engine] Template Path: {PPT_TEMPLATE_PATH} | Exists: {template_exists}")

        if template_exists:
            prs = Presentation(PPT_TEMPLATE_PATH)
        else:
            print("[WARN PPT Engine] 模板文件未找到，降级使用默认空白 Presentation。")
            prs = Presentation()

        cover_layout, content_layout = get_layout_by_type(prs)

        lesson_title = final_basic_info.get("lesson_title") or word_data.get("lesson_title", "教学课件")

        # ---------------- Slide 1: 封面页 ----------------
        slide_cover = prs.slides.add_slide(cover_layout)
        title_ph, subtitle_ph, _ = find_placeholders(slide_cover)

        cover_title_text = f"{lesson_title}" if "课件" in lesson_title else f"{lesson_title} 教学课件"
        if title_ph and title_ph.has_text_frame:
            fill_text_frame(title_ph.text_frame, cover_title_text, align=PP_ALIGN.CENTER)
        else:
            t_box = slide_cover.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2))
            fill_text_frame(t_box.text_frame, cover_title_text, default_size=40, default_color=(0, 51, 102), align=PP_ALIGN.CENTER)

        teacher_name = final_basic_info.get("teacher_name", "教师")
        teacher_unit = final_basic_info.get("teacher_unit", "")
        unit_str = f"（{teacher_unit}）" if teacher_unit else ""
        
        sub_info = (
            f"学科：{final_basic_info.get('subject', '数学')}  |  年级：{final_basic_info.get('grade', '四年级')}\n"
            f"执教教师：{teacher_name}{unit_str}"
        )
        if subtitle_ph and subtitle_ph.has_text_frame:
            fill_text_frame(subtitle_ph.text_frame, sub_info, align=PP_ALIGN.CENTER)
        else:
            s_box = slide_cover.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(2))
            fill_text_frame(s_box.text_frame, sub_info, default_size=20, default_color=(80, 80, 80), align=PP_ALIGN.CENTER)

        # ---------------- Slide 2: 教材与学情分析 ----------------
        textbook_analysis = final_basic_info.get("textbook_analysis")
        student_analysis = final_basic_info.get("student_analysis")
        
        if textbook_analysis or student_analysis:
            slide_analysis = prs.slides.add_slide(content_layout)
            t_ph, _, b_ph = find_placeholders(slide_analysis)
            
            if t_ph and t_ph.has_text_frame:
                fill_text_frame(t_ph.text_frame, "教材与学情分析")
            else:
                tb_title = slide_analysis.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1))
                fill_text_frame(tb_title.text_frame, "教材与学情分析", default_size=32, default_color=(0, 51, 102))

            analysis_text = [
                f"【教材分析】\n{textbook_analysis or '无'}",
                f"【学情分析】\n{student_analysis or '无'}"
            ]
            if b_ph and b_ph.has_text_frame:
                fill_text_frame(b_ph.text_frame, "\n\n".join(analysis_text))
            else:
                tb_body = slide_analysis.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
                fill_text_frame(tb_body.text_frame, "\n\n".join(analysis_text), default_size=16, default_color=(51, 51, 51))

        # ---------------- Slide 3: 教学目标与重难点 ----------------
        slide_target = prs.slides.add_slide(content_layout)
        t_ph, _, b_ph = find_placeholders(slide_target)

        if t_ph and t_ph.has_text_frame:
            fill_text_frame(t_ph.text_frame, "教学目标与重难点")
        else:
            tb_title = slide_target.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1))
            fill_text_frame(tb_title.text_frame, "教学目标与重难点", default_size=32, default_color=(0, 51, 102))

        # 组装教学目标
        objs = word_data.get("teaching_objectives") or word_data.get("targets", {})
        if isinstance(objs, dict):
            k_obj = objs.get("knowledge", "")
            a_obj = objs.get("ability", "")
            l_obj = objs.get("literacy", "")
        else:
            k_obj, a_obj, l_obj = str(objs), "", ""

        key_pts = word_data.get("key_points", "掌握核心概念与表达方法。")
        diff_pts = word_data.get("difficult_points", "理解知识的迁移与运用。")

        target_text = [
            f"【教学重点】\n{key_pts}",
            f"【教学难点】\n{diff_pts}",
            "【核心目标】",
            f"• 知识目标：{k_obj}",
            f"• 能力目标：{a_obj}",
            f"• 素养目标：{l_obj}"
        ]

        if b_ph and b_ph.has_text_frame:
            fill_text_frame(b_ph.text_frame, "\n\n".join(target_text))
        else:
            tb_body = slide_target.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
            fill_text_frame(tb_body.text_frame, "\n\n".join(target_text), default_size=16, default_color=(51, 51, 51))

        # ---------------- Slide 4+: 动态读取教学环节 ----------------
        process_list = word_data.get("teaching_process") or word_data.get("slides") or []

        for idx, item in enumerate(process_list, 1):
            if not isinstance(item, dict):
                continue

            slide_p = prs.slides.add_slide(content_layout)
            t_ph, _, b_ph = find_placeholders(slide_p)

            stage_title = item.get("stage") or item.get("title") or f"环节 {idx}"
            full_title = f"教学环节：{stage_title}" if not stage_title.startswith("教学环节") else stage_title
            
            if t_ph and t_ph.has_text_frame:
                fill_text_frame(t_ph.text_frame, full_title)
            else:
                tb_p_title = slide_p.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1))
                fill_text_frame(tb_p_title.text_frame, full_title, default_size=28, default_color=(0, 51, 102))

            # 正文内容提取
            t_act = item.get("teacher_activity", "")
            s_act = item.get("student_activity", "")
            d_intent = item.get("design_intent", "")
            
            content_text = f"【教师活动】\n{t_act}\n\n【学生活动】\n{s_act}\n\n【设计意图】\n{d_intent}"

            if b_ph and b_ph.has_text_frame:
                fill_text_frame(b_ph.text_frame, content_text)
            else:
                tb_p_body = slide_p.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
                fill_text_frame(tb_p_body.text_frame, content_text, default_size=15, default_color=(51, 51, 51))

        # ---------------- Slide Next: 板书设计与课后作业 ----------------
        board_design = word_data.get("board_design")
        homework = word_data.get("homework")

        if board_design or homework:
            slide_extra = prs.slides.add_slide(content_layout)
            t_ph, _, b_ph = find_placeholders(slide_extra)

            if t_ph and t_ph.has_text_frame:
                fill_text_frame(t_ph.text_frame, "板书设计与作业布置")
            else:
                tb_e_title = slide_extra.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1))
                fill_text_frame(tb_e_title.text_frame, "板书设计与作业布置", default_size=28, default_color=(0, 51, 102))

            extra_lines = []
            if board_design:
                extra_lines.append(f"【板书设计】\n{board_design}")

            if isinstance(homework, dict):
                hw_basic = homework.get("basic", "")
                hw_adv = homework.get("advanced", "")
                extra_lines.append(f"【基础作业】\n{hw_basic}\n\n【拓展作业】\n{hw_adv}")
            elif homework:
                extra_lines.append(f"【课后作业】\n{homework}")

            extra_text = "\n\n".join(extra_lines)

            if b_ph and b_ph.has_text_frame:
                fill_text_frame(b_ph.text_frame, extra_text)
            else:
                tb_e_body = slide_extra.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
                fill_text_frame(tb_e_body.text_frame, extra_text, default_size=15, default_color=(51, 51, 51))

        # 3. 保存输出
        prs.save(output_path)
        print(f"[SUCCESS] PPT 成功生成至：{output_path}")
        return output_path

    except Exception as e:
        print(f"[ERROR PPT Generation] {str(e)}")
        traceback.print_exc()

        # 保底生成
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "教学设计 PPT"
        prs.save(output_path)
        return output_path
# ---------------------------------------------------------------------------

import os
import json
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH

def safe_json_parse(data):
    """确保数据转为 dict 格式"""
    if isinstance(data, dict):
        return data
    if isinstance(data, str) and data.strip():
        try:
            return json.loads(data)
        except Exception:
            return {}
    return {}

def generate_word_lesson_plan(word_data, basic_info, output_path):
    try:
        # 1. 动态安全解析 JSON 数据
        basic_info = safe_json_parse(basic_info)
        word_data = safe_json_parse(word_data)

        # 2. 读取 Word 模板
        if os.path.exists(WORD_TEMPLATE_PATH):
            doc = Document(WORD_TEMPLATE_PATH)
        else:
            doc = Document()

        lesson_title = basic_info.get('lesson_title') or word_data.get('lesson_title', '')

        # A. 替换主标题
        for p in doc.paragraphs:
            if "教学设计" in p.text:
                p.text = f"《{lesson_title}》教学设计" if lesson_title else p.text
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                break

        # B. 填充主表格
        if len(doc.tables) > 0:
            table = doc.tables[0]

            for idx, row in enumerate(table.rows):
                row_text = "".join([c.text.strip() for c in row.cells])

                # 1. 基础信息行（完全按输入动态填写）
                if any(k in row_text for k in ["授课教师", "学段", "授课时间"]):
                    for c_idx, cell in enumerate(row.cells):
                        cell_txt = cell.text.strip()
                        if cell_txt == "授课教师" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("teacher_name", ""))
                        elif cell_txt == "授课教师单位" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("teacher_unit", ""))
                        elif cell_txt == "授课日期" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("lesson_date", ""))
                        elif cell_txt == "学段" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("stage", ""))
                        elif cell_txt == "学科" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("subject", ""))
                        elif cell_txt == "适用年级" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("grade", ""))
                        elif cell_txt == "授课时间" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("lesson_time", ""))
                        elif cell_txt == "课型" and c_idx + 1 < len(row.cells):
                            row.cells[c_idx + 1].text = str(basic_info.get("lesson_type", ""))

                # 2. 文本区域填充
                elif row.cells[0].text.strip().startswith("课题"):
                    row.cells[-1].text = lesson_title

                elif "教材分析" in row.cells[0].text:
                    row.cells[-1].text = str(basic_info.get("textbook_analysis") or word_data.get("textbook_analysis", ""))

                elif "学情分析" in row.cells[0].text:
                    row.cells[-1].text = str(basic_info.get("student_analysis") or word_data.get("student_analysis", ""))

                elif "教学目标" in row.cells[0].text:
                    obj = word_data.get("teaching_objectives") or basic_info.get("teaching_objectives", {})
                    if isinstance(obj, dict):
                        k = obj.get('knowledge') or obj.get('knowledge_and_skills', '')
                        a = obj.get('ability') or obj.get('process_and_methods', '')
                        l = obj.get('literacy') or obj.get('emotions_and_values', '')
                        row.cells[-1].text = f"1. 知识与技能：{k}\n2. 过程与方法：{a}\n3. 情感态度与价值观：{l}"
                    else:
                        row.cells[-1].text = str(obj)

                elif "教学重难点" in row.cells[0].text:
                    kp = word_data.get('key_points', '')
                    dp = word_data.get('difficult_points', '')
                    row.cells[-1].text = f"教学重点：{kp}\n教学难点：{dp}"

                elif "教法与学法" in row.cells[0].text:
                    tm = word_data.get("teaching_methods", "")
                    if isinstance(tm, dict):
                        row.cells[-1].text = f"教法：{tm.get('teacher_method', '')}\n学法：{tm.get('student_method', '')}"
                    else:
                        row.cells[-1].text = str(tm)

                elif "归纳总结" in row_text or "作业设计" in row_text:
                    # 1. 尝试多途径获取作业数据
                    hw = word_data.get("homework", {})
                    if isinstance(hw, str):
                        hw_basic = hw
                        hw_advanced = ""
                    elif isinstance(hw, dict):
                        hw_basic = hw.get("basic") or hw.get("homework", "")
                        hw_advanced = hw.get("advanced", "")
                    else:
                        hw_basic, hw_advanced = "", ""
                
                    # 组合作业文本
                    homework_text = f"1. 基础作业：{hw_basic}\n2. 拓展作业：{hw_advanced}" if hw_advanced else f"{hw_basic}"
                
                    # 2. 动态提取/生成归纳总结（如果模型没单独给 summary，自动从教学过程第 5 环节提取）
                    sh_dict = word_data.get("summary_and_homework", {})
                    summary_text = ""
                    if isinstance(sh_dict, dict):
                        summary_text = sh_dict.get("summary", "")
                
                    if not summary_text:
                        # 保底方案：从教学过程最后一个环节抽取设计意图作为归纳总结
                        process_list = word_data.get("teaching_process", [])
                        if process_list and len(process_list) >= 5:
                            last_stage = process_list[-1]
                            summary_text = f"引导学生回顾本节课条形统计图的特点（标题、横轴、纵轴、直条），归纳1格代表多个单位的画法。{last_stage.get('design_intent', '')}"
                        else:
                            summary_text = "引导学生自主梳理本节课的核心概念，总结条形统计图的绘制步骤与应用技巧。"
                
                    # 3. 动态提取/生成素养发展（如果模型没单独给，自动从教学目标提取）
                    literacy_text = ""
                    if isinstance(sh_dict, dict):
                        literacy_text = sh_dict.get("literacy_development", "")
                    
                    if not literacy_text:
                        objs = word_data.get("teaching_objectives", {})
                        if isinstance(objs, dict):
                            literacy_text = objs.get("literacy", "") or objs.get("ability", "")
                        if not literacy_text:
                            literacy_text = "通过本土农业数据统计，培养数据分析观念，提升运用数学知识解决真实生活问题的核心素养。"
                
                    # 4. 拼装丰富填入 Word 单元格
                    full_text = (
                        f"【归纳总结】\n{summary_text}\n\n"
                        f"【作业设计】\n{homework_text}\n\n"
                        f"【素养发展】\n{literacy_text}"
                    )
                    row.cells[-1].text = full_text

                elif "板书设计" in row_text and idx + 1 < len(table.rows):
                    table.rows[idx + 1].cells[-1].text = str(word_data.get("board_design", ""))

                elif "课后反思" in row_text and idx + 1 < len(table.rows):
                    table.rows[idx + 1].cells[-1].text = str(word_data.get("reflection", ""))

                # 3. 五环节自动化填充算法（自动匹配 5 个环节）
                else:
                    first_cell = row.cells[0].text.strip()
                    process_list = word_data.get('teaching_process', [])

                    stage_keywords = {
                        "创设": ["创设", "导入", "情境"],
                        "探究": ["探究", "建构", "新知"],
                        "理解": ["理解", "应用", "巩固"],
                        "迁移": ["迁移", "拓展", "延伸"],
                        "创新": ["创新", "提升", "综合"]
                    }

                    if isinstance(process_list, list):
                        for item in process_list:
                            if not isinstance(item, dict): continue
                            p_stage = str(item.get("stage", ""))
                            
                            # 寻找匹配的环节
                            for key, kw_list in stage_keywords.items():
                                if any(kw in first_cell for kw in kw_list) and any(kw in p_stage for kw in kw_list):
                                    if len(row.cells) >= 4:
                                        row.cells[1].text = str(item.get("teacher_activity", ""))
                                        row.cells[2].text = str(item.get("student_activity", ""))
                                        row.cells[3].text = str(item.get("design_intent", ""))
                                    break

        doc.save(output_path)
    except Exception as e:
        print(f"[ERROR Word Generation] {str(e)}")
        doc = Document()
        doc.add_heading(f"《{basic_info.get('lesson_title', '教学设计')}》", 0)
        doc.add_paragraph(str(word_data))
        doc.save(output_path)

    return output_path

# ---------------------------------------------------------------------------
# 3. HTTP 服务 Handler
# ---------------------------------------------------------------------------
class SimpleHandler(BaseHTTPRequestHandler):

    def do_GET(self):
        """处理文件下载"""
        parsed_path = urllib.parse.urlparse(self.path).path
        if parsed_path.startswith('/static/'):
            # 1. 提取文件名并进行 URL 解码 (将 %E5%A4%A7%E7%BA%B2 解转回真实中文)
            raw_filename = parsed_path.replace('/static/', '')
            decoded_filename = urllib.parse.unquote(raw_filename)
            
            file_path = os.path.join(DOWNLOAD_DIR, decoded_filename)
            
            # 调试日志：方便在 Render Logs 查看实际找的是什么文件
            print(f"[DEBUG GET] Request file: {decoded_filename}")
            print(f"[DEBUG GET] Full path: {file_path}")
            print(f"[DEBUG GET] File exists? {os.path.exists(file_path)}")

            if os.path.exists(file_path) and os.path.isfile(file_path):
                self.send_response(200)
                mime_type, _ = mimetypes.guess_type(file_path)
                self.send_header('Content-Type', mime_type or 'application/octet-stream')
                # 兼容中文下载名的 Header 标准写法
                self.send_header('Content-Disposition', f"attachment; filename*=UTF-8''{urllib.parse.quote(decoded_filename)}")
                self.send_header('Content-Length', str(os.path.getsize(file_path)))
                self.send_header('Access-Control-Allow-Origin', '*')
                self.end_headers()
                with open(file_path, 'rb') as f:
                    self.wfile.write(f.read())
                return

        # 找不到文件时返回 404
        self.send_response(404)
        self.send_header('Content-Type', 'text/plain; charset=utf-8')
        self.end_headers()
        self.wfile.write(b"File Not Found")

    def do_POST(self):
        """接收 Dify 请求并生成文档"""
        try:
            content_length = int(self.headers.get('Content-Length', 0))
            post_data = self.rfile.read(content_length)
            req_json = json.loads(post_data.decode('utf-8'))

            model_output = req_json.get('model_output', {})
            if isinstance(model_output, str):
                clean_str = model_output.replace("```json", "").replace("```", "").strip()
                model_output = json.loads(clean_str)

            basic_info = model_output.get('basic_info', {})
            ppt_data = model_output.get('ppt_data', {})
            word_data = model_output.get('word_data', {})

            # 过滤掉文件名中的非法字符 (如 《》 / \ : * ? " < > |)
            raw_title = req_json.get('lesson_title') or basic_info.get('lesson_title') or '教学设计'
            clean_title = "".join([c for c in raw_title if c not in r'\/:*?"<>|《》'])

            ppt_filename = f"{clean_title}_教学课件.pptx"
            word_filename = f"{clean_title}_教学设计.docx"
            
            ppt_path = os.path.join(DOWNLOAD_DIR, ppt_filename)
            word_path = os.path.join(DOWNLOAD_DIR, word_filename)

            # 生成文件
            generate_ppt_from_template(ppt_data, basic_info, ppt_path)
            generate_word_lesson_plan(word_data, basic_info, word_path)

            print(f"[DEBUG POST] Generated Word at: {word_path}, exists: {os.path.exists(word_path)}")
            print(f"[DEBUG POST] Generated PPT at: {ppt_path}, exists: {os.path.exists(ppt_path)}")

            domain = "https://dify-file-generator.onrender.com"
            response_data = {
                "status": "success",
                "ppt_url": f"{domain}/static/{urllib.parse.quote(ppt_filename)}",
                "word_url": f"{domain}/static/{urllib.parse.quote(word_filename)}"
            }
            
            self.send_response(200)
            self.send_header('Content-Type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            self.wfile.write(json.dumps(response_data).encode('utf-8'))

        except Exception as e:
            print(f"[ERROR POST] {str(e)}")
            self.send_response(500)
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({"status": "error", "message": str(e)}).encode('utf-8'))

# ---------------------------------------------------------------------------
# 4. 启动 Server 逻辑
# ---------------------------------------------------------------------------
def run(server_class=HTTPServer, handler_class=SimpleHandler, port=10000):
    server_address = ('', port)
    httpd = server_class(server_address, handler_class)
    print(f"Server starting on port {port}...")
    httpd.serve_forever()

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 10000))
    run(port=port)
